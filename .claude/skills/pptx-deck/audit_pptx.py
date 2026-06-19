#!/usr/bin/env python3
"""
audit_pptx.py — Reusable static audit script for native pptxgenjs decks.

Usage:
    python3 scripts/audit_pptx.py <pptx_file> [--forbidden WORD ...] [--expect STRING ...]

Checks:
  1. Editability: count of embedded images in ppt/media
  2. Boundary (OOB): every shape fits within 13.333" x 7.5"
  3. Text overflow: chars-per-line heuristic (flags <3 cpl with >3 chars)
  4. Forbidden words (configurable, default: generic security terms)
  5. Content verification (optional --expect strings)

Exits 0 on PASS, 1 on any FAIL.
"""

import sys, os, zipfile, argparse
from pptx import Presentation

SLIDE_W = 13.333
SLIDE_H = 7.5
DEFAULT_FORBIDDEN = [
    "password", "passwd", "secret", "api_key", "apikey",
    "token", "credential", ".env", "BEGIN PRIVATE KEY",
]


def emu_to_in(v):
    return v / 914400


def check_editability(path):
    count = 0
    try:
        with zipfile.ZipFile(path) as z:
            media = [n for n in z.namelist() if n.startswith("ppt/media/")]
            count = len(media)
    except Exception:
        pass
    ok = count <= 2
    msg = f"  Editability: {count} image(s) in ppt/media ({'PASS' if ok else 'FAIL'})"
    return ok, msg


def check_slides(path):
    prs = Presentation(path)
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        shapes = []
        for shape in slide.shapes:
            left = emu_to_in(shape.left) if hasattr(shape, "left") else 0
            top = emu_to_in(shape.top) if hasattr(shape, "top") else 0
            w = emu_to_in(shape.width) if hasattr(shape, "width") else 0
            h = emu_to_in(shape.height) if hasattr(shape, "height") else 0
            right = left + w
            bottom = top + h
            text = ""
            if shape.has_text_frame:
                text = shape.text_frame.text
            shapes.append({
                "left": left, "top": top, "right": right, "bottom": bottom,
                "w": w, "h": h, "text": text,
            })
        slides.append(shapes)
    return slides, len(prs.slides)


def check_boundary(slides):
    issues = []
    for i, shapes in enumerate(slides, 1):
        for j, s in enumerate(shapes):
            if s["right"] > SLIDE_W + 0.02 and s["w"] > 0.01:
                issues.append(f"  Slide {i} shape {j}: right edge {s['right']:.3f}\" exceeds {SLIDE_W}\" ({s['text'][:40]})")
            if s["bottom"] > SLIDE_H + 0.02 and s["h"] > 0.01:
                issues.append(f"  Slide {i} shape {j}: bottom edge {s['bottom']:.3f}\" exceeds {SLIDE_H}\" ({s['text'][:40]})")
    ok = len(issues) == 0
    msg = f"  Boundary (OOB): {len(issues)} overflow(s) ({'PASS' if ok else 'FAIL'})"
    for iss in issues:
        msg += "\n" + iss
    return ok, msg


def check_text_overflow(slides):
    issues = []
    for i, shapes in enumerate(slides, 1):
        for j, s in enumerate(shapes):
            text = s["text"]
            text_clean = text.replace("\n", " ").strip()
            if not text_clean or s["w"] < 0.01:
                continue
            n_chars = len(text_clean)
            cpi = n_chars / max(s["w"], 0.01)
            if s["w"] < 0.6 and n_chars > 3:
                issues.append(f"  Slide {i} shape {j}: w={s['w']:.2f}\" with {n_chars} chars -> high risk of vertical wrap ({text_clean[:50]})")
            elif cpi > 10 and s["w"] < 1.0 and n_chars > 4:
                issues.append(f"  Slide {i} shape {j}: ~{cpi:.0f} chars/in in {s['w']:.2f}\" wide box -> possible wrap ({text_clean[:50]})")
    ok = len(issues) == 0
    msg = f"  Text overflow heuristic: {len(issues)} warning(s) ({'PASS' if ok else 'FLAG'})"
    for iss in issues:
        msg += "\n" + iss
    return ok, msg


def check_forbidden(slides, forbidden_list):
    all_text = ""
    for shapes in slides:
        for s in shapes:
            all_text += s["text"] + "\n"
    lower = all_text.lower()
    found = []
    for word in forbidden_list:
        if word.lower() in lower:
            found.append(word)
    ok = len(found) == 0
    msg = f"  Forbidden words: {len(found)} found ({'PASS' if ok else 'FAIL'})"
    for w in found:
        msg += f"\n    '{w}' found in text"
    return ok, msg


def check_content(slides, expected):
    all_text = ""
    for shapes in slides:
        for s in shapes:
            all_text += s["text"] + "\n"
    missing = []
    for e in expected:
        if e not in all_text:
            missing.append(e)
    ok = len(missing) == 0
    msg = f"  Content verification: {len(missing)} missing out of {len(expected)} expected ({'PASS' if ok else 'FAIL'})"
    for m in missing:
        msg += f"\n    '{m}' not found"
    return ok, msg


def main():
    parser = argparse.ArgumentParser(description="Audit native pptxgenjs deck")
    parser.add_argument("pptx", help="Path to the .pptx file")
    parser.add_argument("--forbidden", nargs="*", default=[], help="Additional forbidden words (stackable)")
    parser.add_argument("--expect", nargs="*", default=[], help="Expected content strings")
    args = parser.parse_args()

    if not os.path.exists(args.pptx):
        print(f"FAIL: {args.pptx} not found")
        sys.exit(1)

    try:
        import pptx
    except ImportError:
        print("FAIL: python-pptx is not installed. Install with: pip install python-pptx")
        sys.exit(1)

    path = args.pptx
    forbidden = DEFAULT_FORBIDDEN + args.forbidden
    results = []

    ok, msg = check_editability(path)
    results.append((ok, msg))
    print(msg)

    slides, count = check_slides(path)
    print(f"  Slide count: {count}")
    ok, msg = check_boundary(slides)
    results.append((ok, msg))
    print(msg)

    ok, msg = check_text_overflow(slides)
    results.append((ok, msg))
    print(msg)

    ok, msg = check_forbidden(slides, forbidden)
    results.append((ok, msg))
    print(msg)

    if args.expect:
        ok, msg = check_content(slides, args.expect)
        results.append((ok, msg))
        print(msg)

    all_ok = all(r[0] for r in results)
    print(f"\n{'=' * 50}")
    print(f"VERDICT: {'PASS' if all_ok else 'FAIL'}")
    pass_count = sum(1 for r in results if r[0])
    fail_count = sum(1 for r in results if not r[0])
    print(f"PASS={pass_count} FAIL={fail_count}")
    sys.exit(0 if all_ok else 1)


if __name__ == "__main__":
    main()
